import datetime
from zoneinfo import ZoneInfo
import tempfile
import os
from google.adk.agents import Agent, LlmAgent

# Import PowerPoint libraries
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

# Weather and time functions (tools for the agent)
def get_weather(city: str) -> dict:
    """Retrieves the current weather report for a specified city.
    Args:
        city (str): The name of the city for which to retrieve the weather report.
    Returns:
        dict: status and result or error msg.
    """
    if city.lower() == "new york":
        return {
            "status": "success",
            "report": (
                "The weather in New York is sunny with a temperature of 25 degrees"
                " Celsius (41 degrees Fahrenheit)."
            ),
        }
    else:
        return {
            "status": "error",
            "error_message": f"Weather information for '{city}' is not available.",
        }

def get_current_time(city: str) -> dict:
    """Returns the current time in a specified city.
    Args:
        city (str): The name of the city for which to retrieve the current time.
    Returns:
        dict: status and result or error msg.
    """
    if city.lower() == "new york":
        tz_identifier = "America/New_York"
    else:
        return {
            "status": "error",
            "error_message": (
                f"Sorry, I don't have timezone information for {city}."
            ),
        }
    tz = ZoneInfo(tz_identifier)
    now = datetime.datetime.now(tz)
    report = (
        f'The current time in {city} is {now.strftime("%Y-%m-%d %H:%M:%S %Z%z")}'
    )
    return {"status": "success", "report": report}

# PowerPoint generation functions as tools for the agent
def create_new_presentation() -> dict:
    """
    Creates a new PowerPoint presentation object.
    
    Returns:
        dict: Status and the temporary file path where the presentation will be saved.
    """
    # Create a temporary file to store the presentation
    temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.pptx')
    temp_file_path = temp_file.name
    temp_file.close()
    
    # Initialize presentation and save it
    prs = Presentation()
    prs.save(temp_file_path)
    
    return {
        "status": "success",
        "message": "New presentation created successfully",
        "presentation_path": temp_file_path
    }

def add_title_slide(presentation_path: str, title: str, subtitle: str = "") -> dict:
    """
    Adds a title slide to the presentation.
    
    Args:
        presentation_path (str): Path to the PowerPoint presentation file.
        title (str): Title for the slide.
        subtitle (str, optional): Subtitle for the slide.
    
    Returns:
        dict: Status and result or error message.
    """
    try:
        prs = Presentation(presentation_path)
        title_slide_layout = prs.slide_layouts[0]  # Title slide layout
        title_slide = prs.slides.add_slide(title_slide_layout)
        title_slide.shapes.title.text = title
        
        if subtitle != "":
            subtitle_shape = title_slide.placeholders[1]  # Subtitle placeholder
            subtitle_shape.text = subtitle
        
        prs.save(presentation_path)
        return {
            "status": "success",
            "message": "Title slide added successfully",
            "slide_index": len(prs.slides) - 1
        }
    except Exception as e:
        return {
            "status": "error",
            "error_message": f"Failed to add title slide: {str(e)}"
        }

def add_content_slide(presentation_path: str, title: str, content: str = "") -> dict:
    """
    Adds a content slide with title and text to the presentation.
    
    Args:
        presentation_path (str): Path to the PowerPoint presentation file.
        title (str): Title for the slide.
        content (str, optional): Text content for the slide.
    
    Returns:
        dict: Status and result or error message.
    """
    try:
        prs = Presentation(presentation_path)
        content_slide_layout = prs.slide_layouts[1]  # Content slide layout
        content_slide = prs.slides.add_slide(content_slide_layout)
        content_slide.shapes.title.text = title
        
        if content != "":
            # Add content as text box
            left = Inches(1)
            top = Inches(2)
            width = Inches(8)
            height = Inches(4)
            
            textbox = content_slide.shapes.add_textbox(left, top, width, height)
            textbox.text_frame.text = content
        
        prs.save(presentation_path)
        return {
            "status": "success",
            "message": "Content slide added successfully",
            "slide_index": len(prs.slides) - 1
        }
    except Exception as e:
        return {
            "status": "error",
            "error_message": f"Failed to add content slide: {str(e)}"
        }

def add_chart_slide(
    presentation_path: str,
    title: str,
    categories: list[str],
    series_names: list[str],
    series_values: list[list[float]],
    chart_title: str = "",
    chart_type: str = "COLUMN_CLUSTERED"
) -> dict:
    """
    Adds a slide with a chart to the presentation.
    
    Args:
        presentation_path (str): Path to the PowerPoint presentation file.
        title (str): Title for the slide.
        categories (list[str]): List of category names.
        series_names (list[str]): List of series names.
        series_values (list[list[float]]): List of lists containing values for each series.
        chart_title (str): Title for the chart.
        chart_type (str): Type of chart (default: "COLUMN_CLUSTERED").
    
    Returns:
        dict: Status and result or error message.
    """
    try:
        prs = Presentation(presentation_path)
        content_slide_layout = prs.slide_layouts[1]  # Content slide layout
        chart_slide = prs.slides.add_slide(content_slide_layout)
        chart_slide.shapes.title.text = title
        
        # Convert chart_type string to XL_CHART_TYPE enum
        chart_type_map = {
            "COLUMN_CLUSTERED": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "BAR_CLUSTERED": XL_CHART_TYPE.BAR_CLUSTERED,
            "LINE": XL_CHART_TYPE.LINE,
            "PIE": XL_CHART_TYPE.PIE,
            "SCATTER": XL_CHART_TYPE.SCATTER
        }
        xl_chart_type = chart_type_map.get(chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)
        
        # Create chart data object
        chart_data = CategoryChartData()
        chart_data.categories = categories
        
        for i, name in enumerate(series_names):
            chart_data.add_series(name, series_values[i])
        
        # Add chart to the slide
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(5)
        
        chart = chart_slide.shapes.add_chart(
            xl_chart_type, left, top, width, height, chart_data
        ).chart
        
        # Set chart title if provided
        if chart_title != "":
            chart.has_title = True
            chart.chart_title.text_frame.text = chart_title
        
        prs.save(presentation_path)
        return {
            "status": "success",
            "message": "Chart slide added successfully",
            "slide_index": len(prs.slides) - 1
        }
    except Exception as e:
        return {
            "status": "error",
            "error_message": f"Failed to add chart slide: {str(e)}"
        }

def save_presentation(presentation_path: str, output_filename: str) -> dict:
    """
    Saves the presentation with a specific filename.
    
    Args:
        presentation_path (str): Path to the temporary PowerPoint presentation file.
        output_filename (str): Desired output filename.
    
    Returns:
        dict: Status and result or error message.
    """
    try:
        # Ensure output filename has .pptx extension
        if not output_filename.endswith('.pptx'):
            output_filename += '.pptx'
        
        # Load the presentation from the temporary path
        prs = Presentation(presentation_path)
        
        # Save with the new filename
        prs.save(output_filename)
        
        # Remove the temporary file
        os.remove(presentation_path)
        
        return {
            "status": "success",
            "message": f"Presentation saved successfully as {output_filename}",
            "final_path": output_filename
        }
    except Exception as e:
        return {
            "status": "error",
            "error_message": f"Failed to save presentation: {str(e)}"
        }

def create_weather_time_ppt(city: str, output_filename: str = "") -> dict:
    """
    Creates a complete weather and time report PowerPoint presentation for the specified city.
    
    Args:
        city (str): Name of the city for which to create the report.
        output_filename (str): Desired output filename. If empty, uses city name.
    
    Returns:
        dict: Status and result or error message.
    """
    try:
        # Set default output filename if not provided
        if output_filename == "":
            output_filename = f"{city}_weather_time_report.pptx"
        
        # Create new presentation
        new_prs_result = create_new_presentation()
        if new_prs_result["status"] != "success":
            return new_prs_result
        
        temp_path = new_prs_result["presentation_path"]
        
        # Add title slide
        today_date = datetime.datetime.now().strftime('%Y-%m-%d')
        add_title_slide(
            temp_path,
            f"Weather and Time Report: {city}",
            f"Generated on {today_date}"
        )
        
        # Get weather information
        weather_response = get_weather(city)
        
        # Add weather slide
        if weather_response["status"] == "success":
            weather_content = weather_response["report"]
        else:
            weather_content = weather_response["error_message"]
        
        add_content_slide(
            temp_path,
            f"Current Weather in {city}",
            weather_content
        )
        
        # Get time information
        time_response = get_current_time(city)
        
        # Add time slide
        if time_response["status"] == "success":
            time_content = time_response["report"]
        else:
            time_content = time_response["error_message"]
        
        add_content_slide(
            temp_path,
            f"Current Time in {city}",
            time_content
        )
        
        # Add chart slide with example weather data
        add_chart_slide(
            temp_path,
            f"Weather Forecast for {city}",
            ['Monday', 'Tuesday', 'Wednesday', 'Thursday', 'Friday'],
            ['Temperature (°C)', 'Precipitation (mm)'],
            [[25, 24, 27, 23, 22], [0, 5, 15, 2, 0]],
            f'5-Day Weather Forecast for {city}',
            "COLUMN_CLUSTERED"
        )
        
        # Save the final presentation
        save_result = save_presentation(temp_path, output_filename)
        
        return {
            "status": "success",
            "message": f"Weather and time report presentation created successfully",
            "output_path": save_result.get("final_path", output_filename)
        }
    except Exception as e:
        return {
            "status": "error",
            "error_message": f"Failed to create weather and time presentation: {str(e)}"
        }

# Create the main agent with all tools
root_agent = Agent(
    name="weather_time_ppt_agent",
    model="gemini-2.0-flash",
    description=(
        "Agent to answer questions about the time and weather in a city "
        "and generate PowerPoint presentations with that information."
    ),
    instruction=(
        "You are a helpful agent who can answer user questions about the time and weather in a city. "
        "You can also create PowerPoint presentations to summarize this information. "
        "When asked to create a presentation, use the available PowerPoint tools to build a comprehensive report."
    ),
    tools=[
        get_weather,
        get_current_time,
        create_new_presentation,
        add_title_slide,
        add_content_slide,
        add_chart_slide,
        save_presentation,
        create_weather_time_ppt
    ],
)
